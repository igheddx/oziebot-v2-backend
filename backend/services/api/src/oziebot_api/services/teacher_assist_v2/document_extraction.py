"""TeacherAssist v2 document extraction and AI-context helpers."""

from __future__ import annotations

import io
import math
import uuid
import zipfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

import fitz
from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.orm import Session

from oziebot_api.config import Settings
from oziebot_api.models.teacher_assist_pacing_guide_supporting_material import (
    TeacherAssistPacingGuideSupportingMaterial,
)
from oziebot_api.models.teacher_assist_v2_document_extraction import (
    TeacherAssistV2DocumentExtraction,
)
from oziebot_api.models.teacher_assist_v2_instructional_package import (
    TeacherAssistV2PlanningSupplementalMaterial,
)
from oziebot_api.models.teacher_assist_v2_student_submission import TeacherAssistV2StudentSubmission
from oziebot_api.services.teacher_assist.ocr_errors import TeacherAssistOCRProviderError
from oziebot_api.services.teacher_assist.ocr_provider import get_teacher_assist_ocr_provider
from oziebot_api.services.teacher_assist.ocr_provider_config import (
    resolve_teacher_assist_ocr_provider_mode,
)
from oziebot_api.services.teacher_assist.storage import open_teacher_assist_stream

DOCUMENT_EXTRACTION_PREVIEW_LIMIT = 280
DOCUMENT_CONTEXT_ITEM_LIMIT = 20000
DOCUMENT_CONTEXT_TOTAL_LIMIT = 80000
DOCUMENT_EXTRACTION_STATUSES = {
    "not_started",
    "queued",
    "processing",
    "completed",
    "failed",
    "unsupported",
    "manual_text_provided",
}

TXT_MIME_TYPES = {"text/plain", "text/markdown"}
DOCX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
IMAGE_MIME_TYPES = {
    "image/jpeg",
    "image/jpg",
    "image/png",
    "image/webp",
    "image/gif",
    "image/tiff",
}


def _now() -> datetime:
    return datetime.now(UTC)


def _normalize_text(value: str | None) -> str:
    return "\n".join(
        line.strip() for line in (value or "").replace("\r", "\n").split("\n") if line.strip()
    ).strip()


def _preview(value: str | None, *, limit: int = DOCUMENT_EXTRACTION_PREVIEW_LIMIT) -> str | None:
    normalized = _normalize_text(value)
    if not normalized:
        return None
    if len(normalized) <= limit:
        return normalized
    return f"{normalized[: limit - 3].rstrip()}..."


def _token_estimate(value: str | None) -> int | None:
    normalized = _normalize_text(value)
    if not normalized:
        return None
    return max(1, math.ceil(len(normalized) / 4))


def _effective_text(record: TeacherAssistV2DocumentExtraction | None) -> str | None:
    if record is None:
        return None
    edited = _normalize_text(record.teacher_edited_text)
    if edited:
        return edited
    return _normalize_text(record.extracted_text)


def _extract_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _extract_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        parts = [page.get_text("text") for page in doc]
    finally:
        doc.close()
    return "\n".join(part for part in parts if part)


def _extract_docx(file_bytes: bytes) -> str:
    parts: list[str] = []
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
        for name in ("word/document.xml",):
            if name not in archive.namelist():
                continue
            root = ElementTree.fromstring(archive.read(name))
            for paragraph in root.findall(".//w:p", namespace):
                texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
                joined = "".join(texts).strip()
                if joined:
                    parts.append(joined)
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and str(text).strip():
                parts.append(str(text).strip())
    return "\n".join(parts)


def _extract_via_ocr(
    *, settings: Settings, file_bytes: bytes, mime_type: str, original_filename: str
) -> tuple[str, str]:
    provider_name = (settings.teacher_assist_ocr_provider or "mock").strip() or "mock"
    if resolve_teacher_assist_ocr_provider_mode(provider_name) != "real":
        raise TeacherAssistOCRProviderError(
            "Real OCR is not enabled for this file type. Add text manually or enable OCR.",
            error_code="provider_disabled",
        )
    provider = get_teacher_assist_ocr_provider(settings)
    result = provider.extract_text(
        artifact_type="teacher_assist_v2",
        mime_type=mime_type,
        original_filename=original_filename,
        file_bytes=file_bytes,
        settings=settings,
    )
    return result.extracted_text, result.provider


def _extract_file_text(
    *,
    settings: Settings,
    file_bytes: bytes,
    mime_type: str,
    original_filename: str,
) -> tuple[str, str | None, str | None]:
    normalized_mime = (mime_type or "application/octet-stream").strip().lower()
    suffix = Path(original_filename).suffix.lower()
    try:
        if normalized_mime in TXT_MIME_TYPES or suffix in {".txt", ".md"}:
            return "completed", _extract_txt(file_bytes), "local_txt"
        if normalized_mime == "application/pdf" or suffix == ".pdf":
            text = _normalize_text(_extract_pdf(file_bytes))
            if text:
                return "completed", text, "local_pdf"
            ocr_text, provider = _extract_via_ocr(
                settings=settings,
                file_bytes=file_bytes,
                mime_type="application/pdf",
                original_filename=original_filename,
            )
            return "completed", ocr_text, provider
        if normalized_mime in DOCX_MIME_TYPES or suffix == ".docx":
            return "completed", _extract_docx(file_bytes), "local_docx"
        if normalized_mime in PPTX_MIME_TYPES or suffix == ".pptx":
            return "completed", _extract_pptx(file_bytes), "local_pptx"
        if normalized_mime in IMAGE_MIME_TYPES or suffix in {
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".webp",
            ".tiff",
        }:
            ocr_text, provider = _extract_via_ocr(
                settings=settings,
                file_bytes=file_bytes,
                mime_type=normalized_mime,
                original_filename=original_filename,
            )
            return "completed", ocr_text, provider
    except TeacherAssistOCRProviderError as exc:
        return "failed", None, f"{exc}"
    except (fitz.FileDataError, fitz.mupdf.FzErrorFormat, zipfile.BadZipFile, ValueError) as exc:
        return "failed", None, str(exc)
    return "unsupported", None, f"Unsupported file type: {original_filename}"


def serialize_document_extraction(
    record: TeacherAssistV2DocumentExtraction | None,
    *,
    include_full_text: bool = False,
) -> dict[str, Any] | None:
    if record is None:
        return None
    effective_text = _effective_text(record)
    payload = {
        "id": str(record.id),
        "status": record.extraction_status,
        "provider": record.provider,
        "error_message": record.error_message,
        "character_count": record.character_count,
        "token_estimate": record.token_estimate,
        "preview": _preview(effective_text) or record.extracted_text_preview,
        "effective_text_excerpt": effective_text[:DOCUMENT_CONTEXT_ITEM_LIMIT]
        if effective_text
        else None,
        "extracted_text_preview": record.extracted_text_preview,
        "teacher_edited_text": record.teacher_edited_text,
        "has_extracted_text": bool(_normalize_text(record.extracted_text)),
        "has_teacher_edited_text": bool(_normalize_text(record.teacher_edited_text)),
        "has_usable_text": bool(effective_text),
        "manual_override": bool(_normalize_text(record.teacher_edited_text)),
        "created_at": record.created_at.isoformat(),
        "updated_at": record.updated_at.isoformat(),
    }
    if include_full_text:
        payload["effective_text"] = effective_text
        payload["extracted_text"] = record.extracted_text
    return payload


def _upsert_extraction_record(
    db: Session,
    *,
    existing: TeacherAssistV2DocumentExtraction | None,
    values: dict[str, Any],
) -> TeacherAssistV2DocumentExtraction:
    if existing is None:
        existing = TeacherAssistV2DocumentExtraction(
            id=uuid.uuid4(),
            extraction_status="not_started",
            extracted_text=None,
            extracted_text_preview=None,
            character_count=None,
            token_estimate=None,
            provider=None,
            teacher_edited_text=None,
            error_message=None,
            created_at=_now(),
            updated_at=_now(),
            **values,
        )
        db.add(existing)
    else:
        for key, value in values.items():
            setattr(existing, key, value)
        existing.updated_at = _now()
    db.flush()
    return existing


def ensure_supporting_material_extraction_record(
    db: Session,
    *,
    row: TeacherAssistPacingGuideSupportingMaterial,
) -> TeacherAssistV2DocumentExtraction:
    existing = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.supporting_material_id == row.id
        )
    ).one_or_none()
    return _upsert_extraction_record(
        db,
        existing=existing,
        values={
            "tenant_id": row.tenant_id,
            "teacher_user_id": row.uploaded_by_user_id,
            "source_type": "supporting_material",
            "supporting_material_id": row.id,
            "planning_supplemental_material_id": None,
            "student_submission_id": None,
            "pacing_guide_id": row.pacing_guide_id,
            "package_id": None,
            "assignment_id": None,
            "platform_school_year_id": row.platform_school_year_id,
            "catalog_state_id": row.catalog_state_id,
            "catalog_district_id": row.catalog_district_id,
            "catalog_school_id": row.catalog_school_id,
            "catalog_grade_id": row.catalog_grade_id,
            "catalog_subject_id": row.catalog_subject_id,
            "week_start": None,
            "week_end": None,
            "student_number": None,
            "storage_key": row.storage_key or "",
            "original_filename": row.original_filename or row.title,
            "mime_type": row.mime_type or "application/octet-stream",
            "file_size": int(row.file_size or 0),
        },
    )


def ensure_planning_supplemental_extraction_record(
    db: Session,
    *,
    row: TeacherAssistV2PlanningSupplementalMaterial,
) -> TeacherAssistV2DocumentExtraction:
    existing = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.planning_supplemental_material_id == row.id
        )
    ).one_or_none()
    return _upsert_extraction_record(
        db,
        existing=existing,
        values={
            "tenant_id": row.tenant_id,
            "teacher_user_id": row.teacher_user_id,
            "source_type": "planning_supplemental",
            "supporting_material_id": None,
            "planning_supplemental_material_id": row.id,
            "student_submission_id": None,
            "pacing_guide_id": None,
            "package_id": row.package_id,
            "assignment_id": None,
            "platform_school_year_id": row.platform_school_year_id,
            "catalog_state_id": row.catalog_state_id,
            "catalog_district_id": row.catalog_district_id,
            "catalog_school_id": row.catalog_school_id,
            "catalog_grade_id": row.catalog_grade_id,
            "catalog_subject_id": None,
            "week_start": row.week_start,
            "week_end": row.week_end,
            "student_number": None,
            "storage_key": row.storage_key or "",
            "original_filename": row.original_filename or row.title,
            "mime_type": row.mime_type or "application/octet-stream",
            "file_size": int(row.file_size or 0),
        },
    )


def ensure_student_submission_extraction_record(
    db: Session,
    *,
    row: TeacherAssistV2StudentSubmission,
) -> TeacherAssistV2DocumentExtraction:
    existing = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.student_submission_id == row.id
        )
    ).one_or_none()
    return _upsert_extraction_record(
        db,
        existing=existing,
        values={
            "tenant_id": row.tenant_id,
            "teacher_user_id": row.teacher_user_id,
            "source_type": "student_submission",
            "supporting_material_id": None,
            "planning_supplemental_material_id": None,
            "student_submission_id": row.id,
            "pacing_guide_id": None,
            "package_id": None,
            "assignment_id": row.assignment_id,
            "platform_school_year_id": row.platform_school_year_id,
            "catalog_state_id": None,
            "catalog_district_id": row.catalog_district_id,
            "catalog_school_id": row.catalog_school_id,
            "catalog_grade_id": row.catalog_grade_id,
            "catalog_subject_id": row.catalog_subject_id,
            "week_start": None,
            "week_end": None,
            "student_number": row.student_number,
            "storage_key": row.file_key,
            "original_filename": row.original_filename,
            "mime_type": row.mime_type,
            "file_size": row.file_size,
        },
    )


def run_document_extraction(
    db: Session,
    *,
    settings: Settings,
    record: TeacherAssistV2DocumentExtraction,
) -> TeacherAssistV2DocumentExtraction:
    if not record.storage_key:
        record.extraction_status = "failed"
        record.error_message = "No stored file is available for extraction."
        record.updated_at = _now()
        db.flush()
        return record
    record.extraction_status = "processing"
    record.error_message = None
    record.updated_at = _now()
    db.flush()
    with open_teacher_assist_stream(settings, storage_key=record.storage_key) as stream:
        file_bytes = stream.read()
    status, text, provider_or_error = _extract_file_text(
        settings=settings,
        file_bytes=file_bytes,
        mime_type=record.mime_type,
        original_filename=record.original_filename,
    )
    record.extraction_status = status
    record.updated_at = _now()
    if status == "completed":
        normalized = _normalize_text(text)
        if not normalized:
            record.extraction_status = "failed"
            record.error_message = "No readable text was extracted from this file."
            record.extracted_text = None
            record.extracted_text_preview = None
            record.character_count = None
            record.token_estimate = None
            record.provider = provider_or_error
        else:
            record.extracted_text = normalized
            record.extracted_text_preview = _preview(normalized)
            record.character_count = len(normalized)
            record.token_estimate = _token_estimate(normalized)
            record.provider = provider_or_error
            record.error_message = None
    else:
        record.extracted_text = None
        record.extracted_text_preview = None
        record.character_count = None
        record.token_estimate = None
        record.provider = None
        record.error_message = provider_or_error
    db.flush()
    return record


def save_student_submission_response_text(
    db: Session,
    *,
    submission: TeacherAssistV2StudentSubmission,
    response_text: str,
) -> TeacherAssistV2DocumentExtraction:
    normalized = _normalize_text(response_text)
    if not normalized:
        raise ValueError({"response_text": "Student response text is required."})
    record = ensure_student_submission_extraction_record(db, row=submission)
    record.teacher_edited_text = normalized
    if not _normalize_text(record.extracted_text):
        record.extraction_status = "manual_text_provided"
        record.provider = "teacher_manual"
    record.updated_at = _now()
    db.flush()
    return record


def load_supporting_material_extractions(
    db: Session, *, material_ids: list[uuid.UUID]
) -> dict[uuid.UUID, TeacherAssistV2DocumentExtraction]:
    if not material_ids:
        return {}
    rows = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.supporting_material_id.in_(material_ids)
        )
    ).all()
    return {
        row.supporting_material_id: row for row in rows if row.supporting_material_id is not None
    }


def load_planning_supplemental_extractions(
    db: Session, *, material_ids: list[uuid.UUID]
) -> dict[uuid.UUID, TeacherAssistV2DocumentExtraction]:
    if not material_ids:
        return {}
    rows = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.planning_supplemental_material_id.in_(material_ids)
        )
    ).all()
    return {
        row.planning_supplemental_material_id: row
        for row in rows
        if row.planning_supplemental_material_id is not None
    }


def load_submission_extractions(
    db: Session, *, submission_ids: list[uuid.UUID]
) -> dict[uuid.UUID, TeacherAssistV2DocumentExtraction]:
    if not submission_ids:
        return {}
    rows = db.scalars(
        select(TeacherAssistV2DocumentExtraction).where(
            TeacherAssistV2DocumentExtraction.student_submission_id.in_(submission_ids)
        )
    ).all()
    return {row.student_submission_id: row for row in rows if row.student_submission_id is not None}


def build_document_prompt_context(
    materials: list[dict[str, Any]],
    *,
    source_label: str,
    total_char_limit: int = DOCUMENT_CONTEXT_TOTAL_LIMIT,
    item_char_limit: int = DOCUMENT_CONTEXT_ITEM_LIMIT,
) -> dict[str, Any]:
    used: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []
    remaining = total_char_limit
    for material in materials:
        extraction = material.get("extraction") or {}
        entry = {
            "id": str(material.get("id")) if material.get("id") is not None else None,
            "title": material.get("title"),
            "original_filename": material.get("original_filename"),
            "material_kind": material.get("material_kind"),
            "resource_type": material.get("resource_type"),
            "source_label": source_label,
            "extraction_status": extraction.get("status"),
        }
        effective_text = _normalize_text(
            extraction.get("effective_text")
            or extraction.get("effective_text_excerpt")
            or extraction.get("teacher_edited_text")
            or extraction.get("extracted_text")
        )
        if not effective_text:
            skipped.append(
                {
                    **entry,
                    "reason": extraction.get("error_message")
                    or "No extracted text was available for this file.",
                }
            )
            continue
        excerpt = effective_text[: min(item_char_limit, remaining)].strip()
        if not excerpt:
            skipped.append(
                {**entry, "reason": "Prompt text budget was exhausted before this file."}
            )
            continue
        used.append(
            {
                **entry,
                "included_characters": len(excerpt),
                "excerpt": excerpt,
            }
        )
        remaining -= len(excerpt)
        if remaining <= 0:
            break
    return {
        "used_documents": used,
        "skipped_documents": skipped,
        "used_count": len(used),
        "skipped_count": len(skipped),
        "remaining_char_budget": max(0, remaining),
    }


def build_ai_readiness_summary(
    *,
    district_materials: list[dict[str, Any]],
    teacher_materials: list[dict[str, Any]],
) -> dict[str, Any]:
    file_materials = [
        row for row in district_materials + teacher_materials if row.get("material_kind") == "file"
    ]
    completed = 0
    failed = 0
    pending = 0
    usable = 0
    for row in file_materials:
        extraction = row.get("extraction") or {}
        status = extraction.get("status") or "not_started"
        if extraction.get("has_usable_text"):
            usable += 1
        if status in {"completed", "manual_text_provided"}:
            completed += 1
        elif status in {"failed", "unsupported"}:
            failed += 1
        else:
            pending += 1
    return {
        "district_resources_loaded": len(district_materials),
        "teacher_materials_loaded": len(teacher_materials),
        "uploaded_file_count": len(file_materials),
        "extracted_text_available_count": usable,
        "files_completed_count": completed,
        "files_failed_count": failed,
        "files_pending_count": pending,
        "continue_with_filenames_only": pending > 0 or failed > 0,
    }
