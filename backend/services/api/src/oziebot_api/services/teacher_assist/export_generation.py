from __future__ import annotations

from datetime import UTC, datetime, timedelta
import io
import json
import uuid

from sqlalchemy import select
from sqlalchemy.engine import Engine
from sqlalchemy.orm import Session, sessionmaker

from oziebot_api.config import Settings
from oziebot_api.models.teacher_assist_export_artifact import TeacherAssistExportArtifact
from oziebot_api.models.teacher_assist_weekly_plan import TeacherAssistWeeklyPlan
from oziebot_api.models.teacher_assist_workflow import TeacherAssistWorkflow
from oziebot_api.services.teacher_assist.activity_events import record_activity_event
from oziebot_api.services.teacher_assist.constants import (
    ARTIFACT_EXPORT_OUTPUT_REF_TYPE,
    ARTIFACT_EXPORT_WORKFLOW_TYPE,
    EXPORT_PROMPT_VERSION,
    validate_teacher_assist_export_artifact_status,
    validate_teacher_assist_workflow_status,
    validate_teacher_assist_workflow_type,
)
from oziebot_api.services.teacher_assist.export_artifacts import (
    _artifact_download_filename,
    _artifact_mime_type,
    normalize_export_request,
    touch_export_artifact_status,
    validate_weekly_plan_export_eligibility,
)
from oziebot_api.services.teacher_assist.export_templates import build_export_preview
from oziebot_api.services.teacher_assist.storage import save_teacher_assist_bytes
from oziebot_api.services.teacher_assist.workflow_service import get_visible_weekly_plan_or_404

EXPORT_LOG_LIMIT = 40


class TeacherAssistExportCancelledError(RuntimeError):
    pass


def _append_export_log(workflow: TeacherAssistWorkflow, *, event: str, message: str) -> None:
    entries = list(workflow.execution_log_json or [])
    entries.append(
        {
            "event": event,
            "message": message,
            "recorded_at": datetime.now(UTC).isoformat(),
        }
    )
    workflow.execution_log_json = entries[-EXPORT_LOG_LIMIT:]


def _clear_workflow_lease(workflow: TeacherAssistWorkflow) -> None:
    workflow.leased_by_worker = None
    workflow.lease_expires_at = None
    workflow.heartbeat_at = None
    workflow.timeout_at = None


def _touch_workflow_heartbeat(
    workflow: TeacherAssistWorkflow,
    *,
    settings: Settings,
    worker_name: str,
    progress_percent: int | None = None,
) -> None:
    now = datetime.now(UTC)
    workflow.leased_by_worker = worker_name
    workflow.heartbeat_at = now
    workflow.lease_expires_at = now + timedelta(seconds=max(1, settings.teacher_assist_worker_lease_seconds))
    if workflow.timeout_at is None:
        workflow.timeout_at = now + timedelta(seconds=max(1, settings.teacher_assist_extraction_timeout_seconds))
    workflow.updated_at = now
    if progress_percent is not None:
        workflow.progress_percent = progress_percent


def _set_workflow_status(
    workflow: TeacherAssistWorkflow,
    *,
    status: str,
    progress_percent: int | None = None,
    error_message: str | None = None,
) -> None:
    now = datetime.now(UTC)
    workflow.status = validate_teacher_assist_workflow_status(status)
    workflow.updated_at = now
    workflow.error_message = error_message
    if progress_percent is not None:
        workflow.progress_percent = progress_percent
    if status == "running" and workflow.started_at is None:
        workflow.started_at = now
    if status in {"completed", "failed", "cancelled"}:
        workflow.completed_at = now
        _clear_workflow_lease(workflow)


def render_pptx_bytes(preview: dict) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for slide in preview.get("slides", []):
        layout = presentation.slide_layouts[1]
        slide_obj = presentation.slides.add_slide(layout)
        slide_obj.shapes.title.text = str(slide.get("title", "Slide"))
        body = slide_obj.placeholders[1].text_frame
        body.clear()
        bullets = slide.get("bullets") or []
        if not bullets:
            body.text = " "
            continue
        body.text = str(bullets[0])
        for bullet in bullets[1:]:
            body.add_paragraph().text = str(bullet)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def render_printable_html_bytes(preview: dict) -> bytes:
    title = preview.get("title") or "TeacherAssist Export"
    sections: list[str] = [f"<h1>{title}</h1>"]
    if preview.get("artifact_kind") == "slides":
        for slide in preview.get("slides") or []:
            sections.append(f"<h2>{slide.get('title', 'Slide')}</h2>")
            sections.append("<ul>")
            for bullet in slide.get("bullets") or []:
                sections.append(f"<li>{bullet}</li>")
            sections.append("</ul>")
    else:
        for question in preview.get("questions") or []:
            sections.append(f"<h2>{question.get('question_text', 'Question')}</h2>")
            if question.get("choices"):
                sections.append("<ul>")
                for choice in question["choices"]:
                    sections.append(f"<li>{choice}</li>")
                sections.append("</ul>")
            if question.get("answer_key"):
                sections.append(f"<p><strong>Answer key:</strong> {question['answer_key']}</p>")
    html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<title>{title}</title></head><body>{''.join(sections)}</body></html>"
    )
    return html.encode("utf-8")


def _render_export_file_bytes(*, preview: dict, export_format: str) -> bytes:
    if export_format == "pptx":
        return render_pptx_bytes(preview)
    if export_format == "json":
        return json.dumps(preview, indent=2, sort_keys=True).encode("utf-8")
    if export_format == "printable_html":
        return render_printable_html_bytes(preview)
    raise ValueError("Unsupported export format")


def create_weekly_plan_export(
    db: Session,
    *,
    tenant_id: uuid.UUID,
    user_id: uuid.UUID,
    weekly_plan_id: uuid.UUID,
    artifact_type: str,
    export_format: str | None = None,
    provider_mode: str = "mock",
    settings: Settings | None = None,
) -> TeacherAssistExportArtifact:
    settings = settings or Settings()
    validate_weekly_plan_export_eligibility(
        db,
        tenant_id=tenant_id,
        user_id=user_id,
        weekly_plan_id=weekly_plan_id,
    )
    normalized_type, normalized_format = normalize_export_request(
        artifact_type=artifact_type,
        export_format=export_format,
    )
    normalized_provider_mode = provider_mode.strip().lower() or "mock"
    if normalized_provider_mode not in {"mock", "real"}:
        raise ValueError("Unsupported export provider mode")
    if normalized_provider_mode == "real" and not settings.teacher_assist_real_provider_enabled:
        raise ValueError("Real export generation is disabled")

    plan = get_visible_weekly_plan_or_404(
        db,
        tenant_id=tenant_id,
        user_id=user_id,
        weekly_plan_id=weekly_plan_id,
    )
    now = datetime.now(UTC)
    snapshot = {
        "weekly_plan_id": str(plan.id),
        "weekly_plan_title": plan.title,
        "artifact_type": normalized_type,
        "export_format": normalized_format,
        "provider_mode": normalized_provider_mode,
        "content_json": plan.content_json,
    }
    workflow = TeacherAssistWorkflow(
        tenant_id=tenant_id,
        user_id=user_id,
        planning_input_draft_id=plan.planning_input_draft_id,
        workflow_type=validate_teacher_assist_workflow_type(ARTIFACT_EXPORT_WORKFLOW_TYPE),
        status=validate_teacher_assist_workflow_status("queued"),
        input_snapshot_json=snapshot,
        output_ref_type=None,
        output_ref_id=None,
        error_message=None,
        progress_percent=0,
        leased_by_worker=None,
        lease_expires_at=None,
        heartbeat_at=None,
        retry_count=0,
        max_retries=max(0, settings.teacher_assist_worker_max_retries),
        timeout_at=None,
        provider_name=normalized_provider_mode,
        provider_model="mock-export-generator" if normalized_provider_mode == "mock" else settings.teacher_assist_real_provider_model,
        prompt_version=EXPORT_PROMPT_VERSION,
        last_error_code=None,
        execution_log_json=[],
        created_at=now,
        started_at=None,
        completed_at=None,
        updated_at=now,
    )
    db.add(workflow)
    db.flush()

    artifact = TeacherAssistExportArtifact(
        tenant_id=tenant_id,
        user_id=user_id,
        source_plan_id=plan.id,
        source_assignment_id=None,
        workflow_id=workflow.id,
        artifact_type=normalized_type,
        artifact_status=validate_teacher_assist_export_artifact_status("queued"),
        title=f"{plan.title} — {normalized_type.replace('_', ' ').title()}",
        export_format=normalized_format,
        storage_key=None,
        preview_json={},
        metadata_json={"provider_mode": normalized_provider_mode},
        provider_name=workflow.provider_name,
        provider_model=workflow.provider_model,
        prompt_version=workflow.prompt_version,
        created_at=now,
        updated_at=now,
    )
    db.add(artifact)
    db.flush()

    record_activity_event(
        db,
        tenant_id=tenant_id,
        user_id=user_id,
        event_type="export_queued",
        event_category="export",
        entity_type="export_artifact",
        entity_id=artifact.id,
        workflow_id=workflow.id,
        summary_text=f"Queued {normalized_type.replace('_', ' ')} export for weekly plan.",
        details_json={
            "artifact_type": normalized_type,
            "export_format": normalized_format,
            "weekly_plan_id": str(plan.id),
        },
    )
    db.flush()
    return artifact


def reclaim_stale_export_workflows(db: Session) -> int:
    now = datetime.now(UTC)
    stale_rows = db.scalars(
        select(TeacherAssistWorkflow).where(
            TeacherAssistWorkflow.workflow_type == ARTIFACT_EXPORT_WORKFLOW_TYPE,
            TeacherAssistWorkflow.status == "running",
            TeacherAssistWorkflow.lease_expires_at.is_not(None),
            TeacherAssistWorkflow.lease_expires_at < now,
        )
    ).all()
    for workflow in stale_rows:
        _mark_export_workflow_for_retry_or_failure(
            db,
            workflow=workflow,
            exc=TimeoutError("TeacherAssist export workflow lease expired"),
            error_code="lease_expired",
        )
    db.flush()
    return len(stale_rows)


def claim_next_export_workflow(
    db: Session,
    *,
    settings: Settings,
    worker_name: str,
    workflow_id: uuid.UUID | None = None,
) -> TeacherAssistWorkflow | None:
    reclaim_stale_export_workflows(db)
    query = (
        select(TeacherAssistWorkflow)
        .where(
            TeacherAssistWorkflow.workflow_type == ARTIFACT_EXPORT_WORKFLOW_TYPE,
            TeacherAssistWorkflow.status == "queued",
        )
        .order_by(TeacherAssistWorkflow.created_at.asc())
    )
    if workflow_id is not None:
        query = query.where(TeacherAssistWorkflow.id == workflow_id)
    workflow = db.scalars(query).first()
    if workflow is None:
        return None

    _set_workflow_status(workflow, status="running", progress_percent=max(workflow.progress_percent, 5))
    _touch_workflow_heartbeat(
        workflow,
        settings=settings,
        worker_name=worker_name,
        progress_percent=max(workflow.progress_percent, 10),
    )
    _append_export_log(workflow, event="export_claimed", message="Export workflow claimed by worker")
    artifact = db.scalars(
        select(TeacherAssistExportArtifact).where(
            TeacherAssistExportArtifact.workflow_id == workflow.id,
        )
    ).one_or_none()
    if artifact is not None:
        touch_export_artifact_status(artifact, status="generating")
    db.flush()
    return workflow


def _mark_export_workflow_for_retry_or_failure(
    db: Session,
    *,
    workflow: TeacherAssistWorkflow,
    exc: Exception,
    error_code: str,
) -> None:
    attempt_number = workflow.retry_count + 1
    workflow.retry_count = attempt_number
    workflow.last_error_code = error_code
    workflow.error_message = str(exc)
    workflow.updated_at = datetime.now(UTC)
    _append_export_log(workflow, event="export_failed", message=str(exc))
    artifact = db.scalars(
        select(TeacherAssistExportArtifact).where(TeacherAssistExportArtifact.workflow_id == workflow.id)
    ).one_or_none()
    if artifact is not None:
        metadata = dict(artifact.metadata_json or {})
        metadata.update({"error_code": error_code, "error_message": str(exc)})
        artifact.metadata_json = metadata
    if attempt_number <= workflow.max_retries:
        _set_workflow_status(workflow, status="queued", progress_percent=min(max(workflow.progress_percent, 5), 95))
        if artifact is not None:
            touch_export_artifact_status(artifact, status="queued")
        return
    _set_workflow_status(
        workflow,
        status="failed",
        progress_percent=min(max(workflow.progress_percent, 5), 95),
        error_message=str(exc),
    )
    if artifact is not None:
        touch_export_artifact_status(artifact, status="failed")


def _persist_export_success(
    db: Session,
    *,
    workflow_id: uuid.UUID,
    settings: Settings,
    worker_name: str,
) -> None:
    workflow = db.scalars(
        select(TeacherAssistWorkflow).where(TeacherAssistWorkflow.id == workflow_id)
    ).one()
    if workflow.status == "cancelled":
        raise TeacherAssistExportCancelledError("Export workflow was cancelled")
    artifact = db.scalars(
        select(TeacherAssistExportArtifact).where(TeacherAssistExportArtifact.workflow_id == workflow.id)
    ).one()
    plan = db.scalars(
        select(TeacherAssistWeeklyPlan).where(TeacherAssistWeeklyPlan.id == artifact.source_plan_id)
    ).one()
    preview = build_export_preview(
        plan_title=plan.title,
        content_json=dict(plan.content_json or {}),
        artifact_type=artifact.artifact_type,
    )
    file_bytes = _render_export_file_bytes(preview=preview, export_format=artifact.export_format)
    stored = save_teacher_assist_bytes(
        settings,
        tenant_id=artifact.tenant_id,
        area="exports",
        original_filename=_artifact_download_filename(artifact),
        contents=file_bytes,
        mime_type=_artifact_mime_type(artifact.export_format),
    )
    artifact.preview_json = preview
    artifact.storage_key = stored.storage_key
    artifact.metadata_json = {
        **dict(artifact.metadata_json or {}),
        "file_size": stored.file_size,
        "mime_type": stored.mime_type,
        "generator": "mock",
    }
    touch_export_artifact_status(artifact, status="ready")
    workflow.output_ref_type = ARTIFACT_EXPORT_OUTPUT_REF_TYPE
    workflow.output_ref_id = artifact.id
    _set_workflow_status(workflow, status="completed", progress_percent=100)
    _append_export_log(workflow, event="export_completed", message="Export artifact generated successfully")
    record_activity_event(
        db,
        tenant_id=workflow.tenant_id,
        user_id=workflow.user_id,
        event_type="export_completed",
        event_category="export",
        entity_type="export_artifact",
        entity_id=artifact.id,
        workflow_id=workflow.id,
        summary_text=f"Completed {artifact.artifact_type.replace('_', ' ')} export.",
        details_json={
            "artifact_type": artifact.artifact_type,
            "export_format": artifact.export_format,
            "weekly_plan_id": str(artifact.source_plan_id),
        },
    )
    db.flush()


def _persist_export_failure(
    factory,
    *,
    workflow_id: uuid.UUID,
    exc: Exception,
    error_code: str,
) -> None:
    session = factory()
    try:
        workflow = session.scalars(
            select(TeacherAssistWorkflow).where(TeacherAssistWorkflow.id == workflow_id)
        ).one_or_none()
        if workflow is None:
            session.commit()
            return
        _mark_export_workflow_for_retry_or_failure(
            session,
            workflow=workflow,
            exc=exc,
            error_code=error_code,
        )
        session.commit()
    finally:
        session.close()


def process_claimed_export_workflow_with_factory(
    factory,
    workflow_id: uuid.UUID,
    *,
    settings: Settings,
    worker_name: str,
) -> None:
    session = factory()
    try:
        _persist_export_success(
            session,
            workflow_id=workflow_id,
            settings=settings,
            worker_name=worker_name,
        )
        session.commit()
    except TeacherAssistExportCancelledError:
        session.rollback()
    except Exception as exc:
        session.rollback()
        _persist_export_failure(
            factory,
            workflow_id=workflow_id,
            exc=exc,
            error_code="execution_failed",
        )
    finally:
        session.close()


def process_next_teacher_assist_export_with_engine(
    engine: Engine,
    *,
    settings: Settings | None = None,
    worker_name: str = "teacher-assist-worker",
    workflow_id: uuid.UUID | None = None,
) -> uuid.UUID | None:
    settings = settings or Settings()
    factory = sessionmaker(bind=engine, autoflush=False, autocommit=False, expire_on_commit=False)
    claim_session = factory()
    try:
        workflow = claim_next_export_workflow(
            claim_session,
            settings=settings,
            worker_name=worker_name,
            workflow_id=workflow_id,
        )
        if workflow is None:
            claim_session.commit()
            return None
        claimed_id = workflow.id
        claim_session.commit()
    finally:
        claim_session.close()
    process_claimed_export_workflow_with_factory(
        factory,
        claimed_id,
        settings=settings,
        worker_name=worker_name,
    )
    return claimed_id
